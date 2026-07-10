import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors (Modern Clean Light Theme)
BG_COLOR = RGBColor(255, 255, 255)       # White: #FFFFFF
CARD_BG = RGBColor(248, 250, 252)        # Slate-50: #F8FAFC
CARD_BORDER = RGBColor(226, 232, 240)    # Slate-200: #E2E8F0
TEXT_WHITE = RGBColor(15, 23, 42)        # Slate-900: #0F172A (used as header text in cards)
TEXT_LIGHT = RGBColor(51, 65, 85)        # Slate-700: #334155 (main body text)
TEXT_MUTED = RGBColor(100, 116, 139)     # Slate-500: #64748B
ACCENT_BLUE = RGBColor(37, 99, 235)      # Blue-600: #2563EB (headers/titles)
ACCENT_GREEN = RGBColor(5, 150, 105)     # Emerald-600: #059669 (categories/success)
ACCENT_ORANGE = RGBColor(217, 119, 6)    # Amber-600: #D97706

def set_dark_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_slide_header(slide, title, category="PROJET SMART SUPPLY CHAIN"):
    # Category tag
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    cat_tf = cat_box.text_frame
    cat_tf.word_wrap = True
    cat_tf.margin_left = cat_tf.margin_right = cat_tf.margin_top = cat_tf.margin_bottom = 0
    p_cat = cat_tf.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = "Segoe UI"
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_GREEN

    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.6))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.margin_left = title_tf.margin_right = title_tf.margin_top = title_tf.margin_bottom = 0
    p_title = title_tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = ACCENT_BLUE

def add_card(slide, left, top, width, height, title=None, border_color=CARD_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    
    if title:
        # Add textbox over it for title to keep styling control
        title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        return left + Inches(0.2), top + Inches(0.65), width - Inches(0.4), height - Inches(0.8)
        
    return left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4)

def add_bullet_list(slide, left, top, width, height, bullets, text_color=TEXT_LIGHT, font_size=13):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        # Check if bullet has sub-bullets
        if isinstance(bullet, tuple) or isinstance(bullet, list):
            main_text, sub_bullets = bullet
            p.text = "• " + main_text
            p.font.name = "Segoe UI"
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = TEXT_WHITE
            p.space_after = Pt(4)
            p.space_before = Pt(6)
            
            for sub in sub_bullets:
                p_sub = tf.add_paragraph()
                p_sub.text = "    - " + sub
                p_sub.font.name = "Segoe UI"
                p_sub.font.size = Pt(font_size - 1)
                p_sub.font.color.rgb = text_color
                p_sub.space_after = Pt(3)
        else:
            # Normal bullet
            if bullet.startswith("**") and "**:" in bullet:
                parts = bullet.split("**: ", 1)
                bold_part = parts[0].replace("**", "")
                normal_part = parts[1]
                
                p.text = "• "
                p.font.name = "Segoe UI"
                p.font.size = Pt(font_size)
                p.font.color.rgb = text_color
                
                run1 = p.add_run()
                run1.text = bold_part + ": "
                run1.font.name = "Segoe UI"
                run1.font.size = Pt(font_size)
                run1.font.bold = True
                run1.font.color.rgb = ACCENT_BLUE
                
                run2 = p.add_run()
                run2.text = normal_part
                run2.font.name = "Segoe UI"
                run2.font.size = Pt(font_size)
                run2.font.color.rgb = text_color
            else:
                p.text = "• " + bullet
                p.font.name = "Segoe UI"
                p.font.size = Pt(font_size)
                p.font.color.rgb = text_color
            
            p.space_after = Pt(6)
            p.space_before = Pt(3)

def add_framed_image(slide, image_path, left, top, width, height, title=None):
    if not os.path.exists(image_path):
        print(f"Warning: Image {image_path} not found. Adding placeholder card.")
        c_left, c_top, c_width, c_height = add_card(slide, left, top, width, height, title or "Aperçu Visuel")
        txBox = slide.shapes.add_textbox(c_left, c_top, c_width, c_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"\n\n\n[ Image manquante ]\n{os.path.basename(image_path)}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MUTED
        return

    # Draw framing card slightly larger than image
    frame_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    frame_shape.fill.solid()
    frame_shape.fill.fore_color.rgb = CARD_BG
    frame_shape.line.color.rgb = CARD_BORDER
    frame_shape.line.width = Pt(1.5)
    
    # Place image inside frame with a small padding
    padding = Inches(0.1)
    slide.shapes.add_picture(
        image_path, 
        left + padding, 
        top + padding, 
        width - (padding * 2), 
        height - (padding * 2)
    )
    
    if title:
        title_box = slide.shapes.add_textbox(left, top - Inches(0.4), width, Inches(0.35))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title.upper()
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_MUTED

# ==============================================================================
# SLIDE 1: Title Slide (Page de Garde)
# ==============================================================================
slide_layout = prs.slide_layouts[6] # Blank
slide1 = prs.slides.add_slide(slide_layout)
set_dark_background(slide1)

t_left, t_top, t_width, t_height = Inches(1.0), Inches(1.2), Inches(11.333), Inches(4.5)
c_left, c_top, c_width, c_height = add_card(slide1, t_left, t_top, t_width, t_height, border_color=ACCENT_BLUE)

txBox = slide1.shapes.add_textbox(c_left, c_top + Inches(0.2), c_width, Inches(2.2))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "SMART SUPPLY CHAIN DASHBOARD"
p.font.name = "Segoe UI"
p.font.size = Pt(38)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
p.space_after = Pt(8)

p2 = tf.add_paragraph()
p2.text = "Tableau de Bord Logistique Intelligent & Aide à la Décision Prédictive"
p2.font.name = "Segoe UI"
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = TEXT_WHITE
p2.space_after = Pt(15)

p3 = tf.add_paragraph()
p3.text = "Une solution d'aide à la décision intégrant l'IA prédictive pour la planification de la demande,\n" \
          "les pipelines de machine learning pour les opérations logistiques,\n" \
          "et une architecture modulaire pour piloter efficacement les flux physiques de l'entreprise."
p3.font.name = "Segoe UI"
p3.font.size = Pt(13)
p3.font.color.rgb = TEXT_LIGHT

meta_box = slide1.shapes.add_textbox(c_left, c_top + Inches(2.3), c_width, Inches(1.3))
meta_tf = meta_box.text_frame
meta_tf.word_wrap = True

p_meta1 = meta_tf.paragraphs[0]
p_meta1.text = "Auteur : SABOR Abderrahmane  |  Filière : Cybersécurité et Réseaux Intelligents (ENSA Khouribga)"
p_meta1.font.name = "Segoe UI"
p_meta1.font.size = Pt(12)
p_meta1.font.bold = True
p_meta1.font.color.rgb = ACCENT_GREEN
p_meta1.space_after = Pt(4)

p_meta2 = meta_tf.add_paragraph()
p_meta2.text = "Organisme d'accueil : ADDSER CONSEIL  |  Période : Juin - Septembre 2026"
p_meta2.font.name = "Segoe UI"
p_meta2.font.size = Pt(12)
p_meta2.font.bold = True
p_meta2.font.color.rgb = TEXT_WHITE

brand_box = slide1.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.333), Inches(0.4))
brand_tf = brand_box.text_frame
p_brand = brand_tf.paragraphs[0]
p_brand.alignment = PP_ALIGN.RIGHT
p_brand.text = "ADDSER CONSEIL  x  ENSA KHOURIBGA"
p_brand.font.name = "Segoe UI"
p_brand.font.size = Pt(11)
p_brand.font.bold = True
p_brand.font.color.rgb = TEXT_MUTED

# ==============================================================================
# SLIDE 2: Introduction & Contexte
# ==============================================================================
slide2 = prs.slides.add_slide(slide_layout)
set_dark_background(slide2)
add_slide_header(slide2, "Introduction & Problématique Métier")

sc_left, sc_top, sc_width, sc_height = Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1)
c_l, c_t, c_w, c_h = add_card(slide2, sc_left, sc_top, sc_width, sc_height, "Enjeux Supply Chain & Logistique")
sc_bullets = [
    "**Manque de visibilité proactive**: Les décisions de réapprovisionnement se basent sur des historiques statiques plutôt que des analyses prédictives.",
    "**Risque de rupture de stock (Stockouts)**: Entraîne des pertes directes de chiffre d'affaires et une dégradation du taux de service (OTIF).",
    "**Coût de surstockage**: Immobilisation inutile de capital et augmentation des frais d'entreposage logistique.",
    "**Objectif**: Passer d'une logique purement réactive à une planification proactive basée sur des indicateurs d'IA en temps réel."
]
add_bullet_list(slide2, c_l, c_t, c_w, c_h, sc_bullets)

sec_left, sec_top, sec_width, sec_height = Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.1)
c_l2, c_t2, c_w2, c_h2 = add_card(slide2, sec_left, sec_top, sec_width, sec_height, "Objectifs Opérationnels du Projet")
sec_bullets = [
    "**Centralisation des Données**: Rassembler les flux de commandes, d'inventaire, de transport et de fournisseurs dans une interface unifiée.",
    "**Pilotage Prédictif de la Demande**: Intégrer des modèles de prévision de séries temporelles pour anticiper les ventes futures.",
    "**Optimisation Automatisée**: Automatiser le calcul des seuils d'inventaire critique et suggérer les volumes de commandes optimaux.",
    "**Aide à la Décision Naturelle**: Fournir aux gestionnaires un chatbot pour poser des questions complexes sur les données en langage naturel."
]
add_bullet_list(slide2, c_l2, c_t2, c_w2, c_h2, sec_bullets)

# ==============================================================================
# SLIDE 3: Architecture Système & Tech Stack
# ==============================================================================
slide3 = prs.slides.add_slide(slide_layout)
set_dark_background(slide3)
add_slide_header(slide3, "Architecture Système & Technologies")

t_left, t_top, t_width, t_height = Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1)
c_l, c_t, c_w, c_h = add_card(slide3, t_left, t_top, t_width, t_height, "Stack Technologique Moderne")
tech_bullets = [
    ("Frontend (Interface Client)", [
        "Angular 17 : Framework web réactif et performant",
        "TypeScript & RxJS : Gestion propre des flux asynchrones",
        "Chart.js & FullCalendar : Visualisations graphiques interactives"
    ]),
    ("Backend (API REST & Services)", [
        "FastAPI (Python) : Traitement asynchrone ultra-rapide",
        "Uvicorn : Serveur d'exécution réactif",
        "Motor : Driver MongoDB asynchrone non-bloquant"
    ]),
    ("Base de données & Intelligence", [
        "MongoDB : Base de données NoSQL orientée documents",
        "Modèles ML : Prophet (Meta) pour la prévision temporelle",
        "Pipeline ML : LightGBM, KNN, KMeans sous Scikit-Learn",
        "IA & LLM : API Ollama (Qwen2.5 / Llama3 hébergé localement) pour le chatbot"
    ])
]
add_bullet_list(slide3, c_l, c_t, c_w, c_h, tech_bullets, font_size=12)

img_path = "images for mod350/system_structure.png"
add_framed_image(slide3, img_path, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.1), "Schéma de l'Architecture Technique")

# ==============================================================================
# SLIDE 4: Fonctionnalités Principales (MOD350)
# ==============================================================================
slide4 = prs.slides.add_slide(slide_layout)
set_dark_background(slide4)
add_slide_header(slide4, "Fonctionnalités Principales (MOD350)")

f_left, f_top, f_width, f_height = Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1)
c_l, c_t, c_w, c_h = add_card(slide4, f_left, f_top, f_width, f_height, "Tableau de Bord & Gestion Logistique")
feat_bullets = [
    "**Dashboard Métier Unifié**: Affiche les indicateurs logistiques clés : Taux de service (OTIF), Délais de livraison (Lead Times), Chiffre d'Affaires et Alertes.",
    "**Optimisation de l'Inventaire**: Calcul en temps réel et par produit du Stock de Sécurité (SS) et du Point de Commande (ROP) selon la volatilité.",
    "**Recommandations d'Achats**: Détecte automatiquement les produits sous le ROP et génère des suggestions de commandes d'approvisionnement conseillées.",
    "**Gestion des Fournisseurs**: Suivi de la performance de livraison par partenaire et détection proactive des retards."
]
add_bullet_list(slide4, c_l, c_t, c_w, c_h, feat_bullets)

img_path = "images for mod350/AI Supply Chain Dashboard main.png"
add_framed_image(slide4, img_path, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.1), "Écran Principal du Tableau de Bord")

# ==============================================================================
# SLIDE 5: Prédiction de la Demande & Simulation
# ==============================================================================
slide5 = prs.slides.add_slide(slide_layout)
set_dark_background(slide5)
add_slide_header(slide5, "Prédiction de la Demande & Simulation")

df_left, df_top, df_width, df_height = Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1)
c_l, c_t, c_w, c_h = add_card(slide5, df_left, df_top, df_width, df_height, "Modélisation Temporelle & Aide à la Décision")
df_bullets = [
    "**Prévision Temporelle (Meta Prophet)**: Algorithme entraîné sur l'historique des ventes quotidiennes pour projeter la demande future sur 90 jours.",
    "**Stabilité Mathématique**: Application d'une transformation logarithmique sur les volumes pour stabiliser la variance des ventes volatiles.",
    "**Days to Stockout Countdown**: Compte à rebours dynamique par produit ('Jours avant rupture') calculé via stock_actuel / demande_prévue.",
    "**Simulateur Logistique**: Insertion de commandes exceptionnelles simulées pour vérifier si le stock physique résistera au choc de la demande.",
    "**Export CSV**: Téléchargement direct des courbes (historique, prévisions, intervalles de confiance à 95%) pour exploitation ERP."
]
add_bullet_list(slide5, c_l, c_t, c_w, c_h, df_bullets)

img_path = "images for mod350/AI Supply Chain Dashboard demand forecasting.png"
add_framed_image(slide5, img_path, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.1), "Visualisation Interactive des Prévisions")

# ==============================================================================
# SLIDE 6: Ingestion de Données & Pipeline Machine Learning
# ==============================================================================
slide6 = prs.slides.add_slide(slide_layout)
set_dark_background(slide6)
add_slide_header(slide6, "Ingestion Robuste & Pipelines Machine Learning")

ing_left, ing_top, ing_width, ing_height = Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1)
c_l, c_t, c_w, c_h = add_card(slide6, ing_left, ing_top, ing_width, ing_height, "Ingestion & Hot-Swapping Atomique")
ing_bullets = [
    "**Importation Fichiers (CSV/XLSX)**: Système d'importation robuste avec nettoyage automatique des doublons pandas et validation des schémas.",
    "**Entraînement Asynchrone**: Lancement du réentraînement des modèles ML dans des sous-processus OS non-bloquants via asyncio.",
    "**Zero-Downtime Hot-Swapping**: Les anciens poids des modèles restent en mémoire pour répondre aux requêtes de l'API pendant le réentraînement.",
    "**Bascule Atomique**: Une fois le nouveau modèle entraîné, FastAPI remplace le pointeur mémoire par le nouveau modèle de manière instantanée, assurant un service 24/7."
]
add_bullet_list(slide6, c_l, c_t, c_w, c_h, ing_bullets)

ml_left, ml_top, ml_width, ml_height = Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.1)
c_l2, c_t2, c_w2, c_h2 = add_card(slide6, ml_left, ml_top, ml_width, ml_height, "Pipelines de Modélisation IA")
ml_bullets = [
    "**Classification des Ventes (LightGBM)**: Classifie à la volée les caractéristiques de commande pour repérer les anomalies transactionnelles lors de l'importation.",
    "**Détection d'Ecarts Logistiques (KNN)**: Algorithme K-Nearest Neighbors entraîné sur les caractéristiques de livraison pour lever des alertes sur les retards critiques.",
    "**Segmentation Client (KMeans)**: Clustering basé sur la notation RFM (Récence, Fréquence, Montant) pour diviser le portefeuille client en segments stratégiques (Champions, Dormants) et guider l'allocation des stocks en période de tension."
]
add_bullet_list(slide6, c_l2, c_t2, c_w2, c_h2, ml_bullets)

# ==============================================================================
# SLIDE 7: Modélisation de Données & Diagramme de Classes
# ==============================================================================
slide7 = prs.slides.add_slide(slide_layout)
set_dark_background(slide7)
add_slide_header(slide7, "Modélisation de Données & Structure du Code")

class_left, class_top, class_width, class_height = Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1)
c_l, c_t, c_w, c_h = add_card(slide7, class_left, class_top, class_width, class_height, "Structure de Données MongoDB & Objets")
class_bullets = [
    "**MongoDB Document Store**: Données modélisées en documents flexibles pour refléter la structure imbriquée des commandes (Order Lines) et des clients.",
    "**Entités Principales**:",
    ("Modèles de Données", [
        "SalesOrder : Contient l'ID commande, le client, la date, le statut, les délais réels/planifiés et le tableau d'OrderLines.",
        "Product / SKU : Caractéristiques de chaque article (prix, nom, catégorie, taux de discount).",
        "Client / Partenaire : Données client, score RFM calculé et informations d'identification.",
        "Supplier : Informations sur le fournisseur, produits associés et KPIs de livraison."
    ]),
    "**Validation Pydantic**: Validation stricte des types de données à l'entrée de l'API FastAPI pour garantir l'intégrité de la base de données."
]
add_bullet_list(slide7, c_l, c_t, c_w, c_h, class_bullets, font_size=12)

img_path = "class_diagram_latest.png"
add_framed_image(slide7, img_path, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.1), "Diagramme de Classes Technique de l'Application")

# ==============================================================================
# SLIDE 8: Assistant Logistique IA (Chatbot NLP)
# ==============================================================================
slide8 = prs.slides.add_slide(slide_layout)
set_dark_background(slide8)
add_slide_header(slide8, "Assistant Logistique IA (Chatbot NLP)")

chat_left, chat_top, chat_width, chat_height = Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1)
c_l, c_t, c_w, c_h = add_card(slide8, chat_left, chat_top, chat_width, chat_height, "Dialogue en langage naturel avec les données")
chat_bullets = [
    "**Assistant Virtuel IA (NLP)**: Chatbot contextuel (ReAct agent) permettant aux utilisateurs de dialoguer directement avec la base de données de l'application.",
    "**Modèle LLM Local (Ollama)**: Utilise les modèles hébergés localement comme Qwen2.5-7B ou Llama3 pour traduire les phrases de l'utilisateur en requêtes MongoDB.",
    "**Questions Métier Supportées**: Permet de poser des questions logistiques en temps réel :",
    ("Exemples d'utilisation", [
        "« Quel est le niveau de stock actuel du produit SKU 1092 ? »",
        "« Quelles sont les anomalies de livraison actives aujourd'hui ? »",
        "« Donne-moi le taux de service OTIF d'ACME Manufacturing. »"
    ]),
    "**Génération de Synthèses**: Génération automatique de comptes-rendus synthétiques pour l'équipe managériale."
]
add_bullet_list(slide8, c_l, c_t, c_w, c_h, chat_bullets, font_size=12)

img_path = "images for mod350/Ai chatbot pop up.png"
add_framed_image(slide8, img_path, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.1), "Interface du Chatbot IA Intégré")

# ==============================================================================
# SLIDE 9: Synthèse & Perspectives
# ==============================================================================
slide9 = prs.slides.add_slide(slide_layout)
set_dark_background(slide9)
add_slide_header(slide9, "Synthèse et Perspectives d'Évolution")

bil_left, bil_top, bil_width, bil_height = Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1)
c_l, c_t, c_w, c_h = add_card(slide9, bil_left, bil_top, bil_width, bil_height, "Bilan Technique & Métier")
bil_bullets = [
    "**Stack Robuste**: Intégration réussie de technologies asynchrones hautement réactives (**Angular 17 / FastAPI / MongoDB**).",
    "**Aide Décisionnelle Optimale**: Prédictions de vente fiables sur 90 jours avec Prophet et simulation de commandes en direct pour guider les réapprovisionnements.",
    "**Dialogue Innovant**: Interface naturelle fluide avec le chatbot IA connecté à la base de données."
]
add_bullet_list(slide9, c_l, c_t, c_w, c_h, bil_bullets)

per_left, per_top, per_width, per_height = Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.1)
c_l2, c_t2, c_w2, c_h2 = add_card(slide9, per_left, per_top, per_width, per_height, "Perspectives d'Évolution futures")
per_bullets = [
    "**Prévisions Logistiques Multivariées**: Intégrer des données exogènes comme les perturbations météo mondiales, le calendrier des vacances ou les indices de fret pour affiner la prévision Prophet.",
    "**Modélisation d'Optimisation de Tournées**: Coupler le dashboard avec des solveurs de recherche opérationnelle pour suggérer des itinéraires de transport optimaux.",
    "**Déploiement Cloud & Monitoring**: Migration vers un cluster de conteneurs avec monitoring des dérives de modèles (Model Drift) et réentraînement planifié périodiquement."
]
add_bullet_list(slide9, c_l2, c_t2, c_w2, c_h2, per_bullets)


# Save presentation in outputs/
output_path = "outputs/Smart_Supply_Chain_Dashboard_Presentation.pptx"
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"Presentation saved successfully to: {output_path}")

# Save presentation in output/
output_path_single = "output/Smart_Supply_Chain_Dashboard_Presentation.pptx"
os.makedirs(os.path.dirname(output_path_single), exist_ok=True)
prs.save(output_path_single)
print(f"Presentation saved successfully to: {output_path_single}")
